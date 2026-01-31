from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "넷플릭스 락인(Lock-in) 전략 보고서"
    subtitle.text = "흑백요리사 IP를 활용한 유저 리텐션 및 수익화 전략\n\nOTT 시장 분석 프로젝트팀"

    # --- Slide 2: Status ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "1. 현황 및 성과: 시장권을 탈환한 구원 투수"
    
    content = slide.placeholders[1]
    content.text = "📊 주요 성과 지표\n" \
                   "- 시장 점유율 35% 달성 (업계 1위 수성)\n" \
                   "- 브랜드 지수 V자 반등 (방영 전 8위 -> 방영 후 1위)\n" \
                   "- 시즌2 참여 밀도 1.82배 급증 (조회수 대비 댓글 비율)\n\n" \
                   "💡 핵심 인사이트\n" \
                   "\"조회수는 시즌1이 높았으나, 실질적인 팬덤의 참여(Engagement)는 시즌2가 압도적임.\""

    # --- Slide 3: Diagnosis (Valuation) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "2. 핵심 진단: 셰프 IP 가치 평가"
    
    content = slide.placeholders[1]
    content.text = "💰 Media Value vs Viral Efficiency\n" \
                   "- 미디어 노출 가치 1위: 안성재 (약 35억 원)\n" \
                   "- 바이럴 효율(가성비) 1위: 임성근 (2.21점)\n\n" \
                   "🏆 최강록 Fandom Analysis\n" \
                   "- 댓글 참여율 1위 (0.057%)\n" \
                   "- 밈(Meme) 파워: '조리보이', '나야 들기름', '만화책'\n" \
                   "- 결론: 단순 시청자가 아닌 '행동하는 팬덤'을 보유한 핵심 자산"

    # --- Slide 4: Diagnosis (Risk) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "2. 핵심 진단: 이탈 위기 (Risk Warning)"
    
    content = slide.placeholders[1]
    content.text = "🚨 Churn Warning Signals\n" \
                   "- 이탈 위험군: 4,356명 (14일 이상 미활동)\n" \
                   "- 예상 손실액: 연간 약 2.1억 원 (구독료 기준)\n\n" \
                   "🔍 이탈의 근본 원인 (Root Cause)\n" \
                   "1. 인성/도덕성 논란 (51.5%): 출연자 검증 실패에 대한 실망\n" \
                   "2. 공정성 불만 (23.2%): 인기 셰프(최강록)의 분량 실종\n" \
                   "3. 가격 저항 (15.3%): 불쾌한 경험 대비 높은 구독료 인식"

    # --- Slide 5: Strategy ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "3. 전략 제안: 골든 타임 & IP 유니버스"
    
    content = slide.placeholders[1]
    content.text = "⏰ 1. 골든 타임 공략 (Push Notification)\n" \
                   "- 수요일 밤 10시 (탐색 마비): 최강록 B-컷 공개\n" \
                   "- 토요일 새벽 4시 (밤샘 정주행): 임성근 가성비 레시피 공개\n\n" \
                   "🚀 2. 단계별 로드맵\n" \
                   "- Step 1 (구독 가치): 식당 예약 패스트트랙 제공 (가격 저항 해소)\n" \
                   "- Step 2 (신뢰 회복): 시즌3 블라인드 국민 심사단 도입\n" \
                   "- Step 3 (확장): 최강록 x 임성근 스핀오프 런칭"

    # --- Slide 6: Vision ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "4. 결론: Paradigm Shift"
    
    content = slide.placeholders[1]
    content.text = "🚩 New Paradigm: Engagement Platform\n" \
                   "- View (조회수) → Engagement (참여)\n" \
                   "- Traffic (숫자) → Fans (팬덤)\n" \
                   "- Subscriber (구독자) → Advocate (지지자)\n\n" \
                   "💡 Final Thought\n" \
                   "\"콘텐츠가 플랫폼을 이깁니다. 이제 단순 OTT를 넘어,\n" \
                   "팬덤이 놀고 즐기는 '캐릭터 유니버스'로 진화해야 합니다.\""

    # Save
    prs.save('Netflix_Strategy_Report.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
