from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# PPT 초기화
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 비율
prs.slide_height = Inches(7.5)

# 공통 스타일 설정 함수
def set_slide_background(slide, color_rgb):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb

def add_title_subtitle(slide, title_text, subtitle_text):
    # 제목
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # 레드 라인
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(229, 9, 20)
    line.line.fill.background()
    
    # 서브타이틀
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(0.5))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.text = subtitle_text
    sp.font.size = Pt(16)
    sp.font.italic = True
    sp.font.color.rgb = RGBColor(229, 9, 20)

# 이미지 경로 (사용자 환경에 맞춤)
base_path = r"c:\Users\pc\Desktop\mypyproject\black\흑백요리사"

# --- Slide 1: Title ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
tf = title.text_frame
p = tf.paragraphs[0]
p.text = "넷플릭스 락인(Lock-in)의 실체:\n'흑백요리사' IP가 방어한 2.1억 원의 구독 가치"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# --- Slide 2: Key Insights ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "Executive Summary: Key Insights", "전략적 핵심 요약")
body = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(5))
tf = body.text_frame
insights = [
    "• 진짜 성과는 '참여'다: 댓글 참여 밀도 시즌 2가 1.82배(조회수 대비) 상향 평준화",
    "• 숨겨진 가성비 1위: 임성근 셰프 바이럴 강도 2.21점 기록 (효율적 캐스팅의 표준)",
    "• 2.1억의 손실 리스크: 14일 이상 미활동 위험팬 4,300명 방치 시 구독 누수 발생",
    "• 해지의 본질: 조작 논란보다 '인성 검증 부실'과 '구독 가치 저항'이 핵심 원인"
]
for text in insights:
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.space_after = Pt(20)

# --- Slide 3: Status (Page 0) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "1. 현황 및 성과: 시장권을 탈환한 구원 투수", "넷플릭스 점유율 35% 돌파 및 V자 반등")
img_path = os.path.join(base_path, "report_page_0_performance.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))

# --- Slide 4: Diagnosis (Asset Valuation) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "2. 핵심 진단: IP 경제적 파급 효과", "안성재 35억 vs 백종원 22억 가치 분석")
img_path = os.path.join(base_path, "report_page_2_asset.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))

# --- Slide 5: Choi Kang-rok Case Study ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "최강록: 수치 너머의 '진짜 화력'", "댓글 참여율 1위(0.057%)의 슈퍼 IP 캐릭터 분석")
img_path = os.path.join(base_path, "choi_wordcloud.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))
text_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(5))
tf = text_box.text_frame
for line in ["1위: 조리보이 (Meme Power)", "2위: 바질을 곁들인... (Identity)", "3위: 나야 들기름 (Interactive)", "4위: 만화책 (Origin Story)"]:
    p = tf.add_paragraph()
    p.text = line
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)

# --- Slide 6: Risk Diagnosis (Legacy Risk) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "3. 위기 진단: S1 리스크의 전이", "출연진 인성 논란으로 인한 브랜드 신뢰도 타격")
img_path = os.path.join(base_path, "report_page_1_diagnosis.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))

# --- Slide 7: Churn Risk Analysis ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "4. 위험 분석: 2.1억 원의 구독 누수 위기", "이탈 위험군(At Risk) 4,356명 정밀 진단")
img_path = os.path.join(base_path, "rfm_segment_analysis_ko.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2.2), height=Inches(4.5))

# --- Slide 8: Strategic Roadmap ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "5. 전략 제안: IP 가치 극대화 로드맵", "Defense - Quality - Expansion 3단계 전략")
img_path = os.path.join(base_path, "report_page_3_roadmap.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))

# --- Slide 9: Execution Plan ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
add_title_subtitle(slide, "6. 실행 과제: 2.1억을 지키는 3대 액션", "Golden Time Care (매주 수요일 22:00 PUSH)")
img_path = os.path.join(base_path, "report_page_4_execution.png")
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), height=Inches(5))

# --- Final Slide ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, RGBColor(13, 13, 13))
title = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
p = title.text_frame.paragraphs[0]
p.text = "감사합니다.\n넷플릭스의 새로운 승리 방정식을 만듭니다."
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# 저장
pptx_file = os.path.join(base_path, "넷플릭스_락인_전략_보고서_최종.pptx")
prs.save(pptx_file)
print(f"PPT 보고서가 성공적으로 생성되었습니다: {pptx_file}")
